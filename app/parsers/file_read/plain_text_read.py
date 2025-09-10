from __future__ import annotations

import re
from pathlib import Path
from typing import Iterable, List, Set

import chardet
from bs4 import BeautifulSoup
from docx import Document
import openpyxl
import pdfplumber
import xml.etree.ElementTree as ET
from pptx import Presentation  # type: ignore[import-not-found]


def _normalize_whitespace(text: str) -> str:
    """压缩空白字符，移除多余空行，返回纯净文本。"""
    # 将各种空白字符标准化为单个空格/换行
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # 去除行首尾空白
    lines: List[str] = [re.sub(r"\s+", " ", line).strip() for line in text.split("\n")]
    # 移除连续空行
    normalized_lines: List[str] = []
    for line in lines:
        if line == "":
            if normalized_lines and normalized_lines[-1] == "":
                continue
        normalized_lines.append(line)
    return "\n".join(normalized_lines).strip()


def _read_text_with_encoding_detection(path: Path) -> str:
    raw = path.read_bytes()
    detection = chardet.detect(raw)
    encoding = detection.get("encoding") or "utf-8"
    try:
        return raw.decode(encoding, errors="ignore")
    except LookupError:
        # 回退到 utf-8
        return raw.decode("utf-8", errors="ignore")


def _read_html(file_path: str) -> str:
    html_content = _read_text_with_encoding_detection(Path(file_path))
    soup = BeautifulSoup(html_content, "html.parser")
    # 获取纯文本
    return _normalize_whitespace(soup.get_text(separator="\n"))


def _read_pdf(file_path: str) -> str:
    texts: List[str] = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            if page_text:
                texts.append(page_text)
    return _normalize_whitespace("\n".join(texts))


def _read_docx(file_path: str) -> str:
    doc = Document(file_path)
    blocks: List[str] = []
    # 段落
    for para in doc.paragraphs:
        if para.text:
            blocks.append(para.text)
    # 表格
    for table in doc.tables:
        for row in table.rows:
            cells_text: List[str] = []
            for cell in row.cells:
                cell_text = "\n".join(p.text for p in cell.paragraphs if p.text)
                cells_text.append(cell_text)
            if any(cells_text):
                blocks.append("\t".join(cells_text))
    return _normalize_whitespace("\n".join(blocks))


def _read_xlsx(file_path: str) -> str:
    wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    sheets_text: List[str] = []
    for sheet in wb.worksheets:
        rows_text: List[str] = []
        for row in sheet.iter_rows(values_only=True):
            values: List[str] = []
            for value in row:
                if value is None:
                    values.append("")
                else:
                    values.append(str(value))
            rows_text.append("\t".join(values).rstrip())
        sheet_text = "\n".join(r for r in rows_text if r is not None)
        if sheet_text.strip():
            sheets_text.append(sheet_text)
    return _normalize_whitespace("\n\n".join(sheets_text))


def _read_xml(file_path: str) -> str:
    tree = ET.parse(file_path)
    root = tree.getroot()

    def iter_text_nodes(node) -> Iterable[str]:
        if node.text and node.text.strip():
            yield node.text
        for child in list(node):
            yield from iter_text_nodes(child)
        if node.tail and node.tail.strip():
            yield node.tail

    return _normalize_whitespace("\n".join(iter_text_nodes(root)))


def _read_rtf(file_path: str) -> str:
    # 朴素 RTF 文本提取：
    # 1) 处理十六进制转义 \'hh
    # 2) 去除控制字和分组
    raw = Path(file_path).read_text(encoding="latin1", errors="ignore")

    def _hex_to_char(match: re.Match[str]) -> str:
        try:
            return bytes.fromhex(match.group(1)).decode("latin1")
        except Exception:
            return ""

    text = re.sub(r"\\'([0-9a-fA-F]{2})", _hex_to_char, raw)
    # 去除控制字，例如 \b0, \par, \u-?
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)
    # 去除分组符号 { }
    text = re.sub(r"[{}]", "", text)
    return _normalize_whitespace(text)


def _read_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    texts: List[str] = []
    for slide in prs.slides:
        slide_lines: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs if run.text)
                    if line:
                        slide_lines.append(line)
            if hasattr(shape, "table"):
                table = shape.table
                for row in table.rows:
                    cells_text = []
                    for cell in row.cells:
                        cell_text = "\n".join(
                            run.text
                            for p in cell.text_frame.paragraphs
                            for run in p.runs
                            if run.text
                        ) if hasattr(cell, "text_frame") else cell.text
                        cells_text.append(cell_text)
                    slide_lines.append("\t".join(filter(None, cells_text)))
        if slide_lines:
            texts.append("\n".join(slide_lines))
    return _normalize_whitespace("\n\n".join(texts))


def read_text(file_path: str, suffix: str) -> str:
    """读取各种常见文件为纯文本。

    仅返回无格式文本，不进行 Markdown/富文本渲染。HTML 使用 bs4 提取，
    PDF 使用 pdfplumber，DOCX 使用 python-docx，XLSX 使用 openpyxl。
    其他如 XML、RTF、PPTX 做最小实现的文本抽取。
    """
    path = Path(file_path)
    if not path.exists() or not path.is_file():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    normalized_suffix = (suffix or path.suffix).lower()

    # 直接按文本方式读取的后缀
    direct_text_suffixes: Set[str] = {
        ".txt",
        ".log",
        ".csv",
        ".tsv",
        ".json",
        ".ndjson",
        ".yaml",
        ".yml",
        ".ini",
        ".cfg",
        ".conf",
        ".properties",
        ".toml",
        ".sql",
        ".md",
        ".rst",
        ".py",
        ".java",
        ".js",
        ".ts",
        ".css",
        ".scss",
        ".less",
        ".env",
        ".bat",
        ".sh",
        ".go",
        ".rs",
        ".c",
        ".cpp",
        ".h",
        ".hpp",
    }

    if normalized_suffix in direct_text_suffixes:
        return _normalize_whitespace(_read_text_with_encoding_detection(path))

    # HTML
    if normalized_suffix in {".html", ".htm"}:
        return _read_html(file_path)

    # PDF
    if normalized_suffix == ".pdf":
        return _read_pdf(file_path)

    # Word (已由上游将 .doc 转换为 .docx)
    if normalized_suffix == ".docx":
        return _read_docx(file_path)

    # Excel（.xls 建议由上游转为 .xlsx）
    if normalized_suffix == ".xlsx":
        return _read_xlsx(file_path)

    # XML
    if normalized_suffix in {".xml"}:
        return _read_xml(file_path)

    # RTF（朴素实现）
    if normalized_suffix == ".rtf":
        return _read_rtf(file_path)

    # PPTX（若未安装依赖，调用方应在环境中安装 python-pptx）
    if normalized_suffix == ".pptx":
        return _read_pptx(file_path)

    # 音频文件 - 直接调用音频读取器
    audio_suffixes = {".mp3", ".wav", ".flac", ".mp4", ".m4a"}
    if normalized_suffix in audio_suffixes:
        try:
            from app.parsers.file_read.audio_read import read_audio
            return read_audio(file_path, normalized_suffix)
        except Exception as e:
            raise ValueError(f"音频文件处理失败: {str(e)}")

    raise ValueError(f"不支持的文件类型（plain_text）：{normalized_suffix}")

